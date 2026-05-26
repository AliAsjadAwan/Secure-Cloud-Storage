"""
Generate PPTX Presentation for Secure Cloud Storage System
Information Security Project
Authors: Ali Asjad Awan (FA23-BCS-028), Muhammad Sohaib Liaqat (FA23-BCS-130)

Modern dark theme with purple accents, morph transitions, proper sizing.
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn, nsmap
from lxml import etree
import copy
import os

# ═══════════════════════════════════════════════════════════════
# CONSTANTS
# ═══════════════════════════════════════════════════════════════
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# Colors
BG_DARK = RGBColor(0x08, 0x0B, 0x1E)
BG_CARD = RGBColor(0x10, 0x13, 0x30)
BG_CARD_LIGHT = RGBColor(0x16, 0x1A, 0x3E)
PURPLE = RGBColor(0xA8, 0x55, 0xF7)
PURPLE_LIGHT = RGBColor(0xC0, 0x84, 0xFC)
PURPLE_DARK = RGBColor(0x7E, 0x22, 0xCE)
WHITE = RGBColor(0xF0, 0xF0, 0xFF)
TEXT_SEC = RGBColor(0xA0, 0xA4, 0xC0)
TEXT_DIM = RGBColor(0x70, 0x74, 0x90)
GREEN = RGBColor(0x4A, 0xDE, 0x80)
RED = RGBColor(0xF8, 0x71, 0x71)
YELLOW = RGBColor(0xFB, 0xBF, 0x24)
TABLE_HEADER_BG = RGBColor(0x1E, 0x16, 0x48)
TABLE_ROW_BG = RGBColor(0x0D, 0x10, 0x28)
TABLE_ROW_ALT = RGBColor(0x11, 0x14, 0x30)

# Layout margins
MARGIN_LEFT = Inches(0.7)
MARGIN_TOP = Inches(0.4)
CONTENT_TOP = Inches(1.55)
CONTENT_W = Inches(11.9)
TITLE_W = Inches(11.9)
TITLE_H = Inches(0.85)


def add_morph_transition(slide):
    """Add morph transition to a slide via XML."""
    try:
        transition = etree.SubElement(
            slide._element,
            qn('p:transition'),
            attrib={'spd': 'med', 'advClick': '1'}
        )
        # Register p159 namespace for morph  
        nsmap_uri = 'http://schemas.microsoft.com/office/powerpoint/2022/main'
        morph = etree.SubElement(
            transition,
            '{http://schemas.microsoft.com/office/powerpoint/2015/main}morph',
            attrib={'option': 'byObject'}
        )
    except Exception:
        # Fallback: add a fade transition
        try:
            transition = etree.SubElement(
                slide._element,
                qn('p:transition'),
                attrib={'spd': 'med', 'advClick': '1'}
            )
            etree.SubElement(transition, qn('p:fade'))
        except Exception:
            pass


def add_fade_transition(slide):
    """Add fade transition to a slide."""
    try:
        transition = etree.SubElement(
            slide._element,
            qn('p:transition'),
            attrib={'spd': 'med', 'advClick': '1'}
        )
        etree.SubElement(transition, qn('p:fade'))
    except Exception:
        pass


def set_slide_bg(slide, color):
    """Set solid background color for a slide."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape_bg(slide, left, top, width, height, fill_color, border_color=None, corner_radius=None):
    """Add a rounded rectangle background shape."""
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1)
    else:
        shape.line.fill.background()
    if corner_radius is not None:
        # Set corner radius via XML
        sp = shape._element
        prstGeom = sp.find(qn('a:prstGeom'), sp.nsmap) or sp.find('.//' + qn('a:prstGeom'))
        if prstGeom is None:
            spPr = sp.find(qn('p:spPr'))
            if spPr is not None:
                prstGeom = spPr.find(qn('a:prstGeom'))
    shape.shadow.inherit = False
    return shape


def add_textbox(slide, left, top, width, height):
    """Add a textbox and return it."""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.word_wrap = True
    return txBox


def add_text_run(tf, text, size=14, color=WHITE, bold=False, italic=False, font_name='Calibri', alignment=PP_ALIGN.LEFT, spacing_after=Pt(4), spacing_before=Pt(0)):
    """Add a paragraph with a single run to a text frame."""
    p = tf.add_paragraph()
    p.alignment = alignment
    p.space_after = spacing_after
    p.space_before = spacing_before
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font_name
    return p


def set_first_para(tf, text, size=14, color=WHITE, bold=False, alignment=PP_ALIGN.LEFT, font_name='Calibri'):
    """Set the default (first) paragraph text."""
    tf.paragraphs[0].alignment = alignment
    tf.paragraphs[0].space_after = Pt(4)
    run = tf.paragraphs[0].add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    run.font.name = font_name


def add_title_bar(slide, title_text, subtitle_text=None):
    """Add a styled title area with accent line."""
    # Title text
    txBox = add_textbox(slide, MARGIN_LEFT, MARGIN_TOP, TITLE_W, TITLE_H)
    tf = txBox.text_frame
    tf.word_wrap = True
    set_first_para(tf, title_text, size=32, color=WHITE, bold=True, font_name='Calibri')
    
    # Accent line
    line = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGIN_LEFT, Inches(1.15), Inches(1.0), Pt(4)
    )
    line.fill.solid()
    line.fill.fore_color.rgb = PURPLE
    line.line.fill.background()
    line.shadow.inherit = False
    
    if subtitle_text:
        sub = add_textbox(slide, MARGIN_LEFT, Inches(1.32), TITLE_W, Inches(0.4))
        set_first_para(sub.text_frame, subtitle_text, size=15, color=TEXT_SEC, font_name='Calibri')


def add_card(slide, left, top, width, height, icon, title, desc, border_color=None):
    """Add a card with icon, title, description."""
    bc = border_color or RGBColor(0x2A, 0x20, 0x55)
    bg = add_shape_bg(slide, left, top, width, height, BG_CARD, bc)
    
    # Icon
    icon_box = add_textbox(slide, left + Inches(0.2), top + Inches(0.15), Inches(0.6), Inches(0.5))
    set_first_para(icon_box.text_frame, icon, size=24, color=PURPLE, font_name='Segoe UI Emoji')
    
    # Title
    title_box = add_textbox(slide, left + Inches(0.2), top + Inches(0.6), width - Inches(0.4), Inches(0.4))
    set_first_para(title_box.text_frame, title, size=14, color=WHITE, bold=True, font_name='Calibri')
    
    # Description
    desc_box = add_textbox(slide, left + Inches(0.2), top + Inches(0.95), width - Inches(0.4), height - Inches(1.1))
    set_first_para(desc_box.text_frame, desc, size=11, color=TEXT_SEC, font_name='Calibri')


def style_table(table, header_bg=TABLE_HEADER_BG, row_bg=TABLE_ROW_BG, alt_bg=TABLE_ROW_ALT):
    """Style a table with dark theme colors."""
    # Disable default banding
    tbl = table._tbl
    tblPr = tbl.find(qn('a:tblPr'))
    if tblPr is None:
        tblPr = etree.SubElement(tbl, qn('a:tblPr'))
    tblPr.set('bandRow', '0')
    tblPr.set('bandCol', '0')
    tblPr.set('firstRow', '0')
    tblPr.set('lastRow', '0')
    
    # Remove default table style
    tblStyle = tblPr.find(qn('a:tblStyle'))
    if tblStyle is not None:
        tblPr.remove(tblStyle)
    
    for row_idx, row in enumerate(table.rows):
        for cell in row.cells:
            # Background
            if row_idx == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = header_bg
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = row_bg if row_idx % 2 == 1 else alt_bg
            
            # Text
            for para in cell.text_frame.paragraphs:
                para.font.size = Pt(11)
                para.font.name = 'Calibri'
                if row_idx == 0:
                    para.font.color.rgb = PURPLE_LIGHT
                    para.font.bold = True
                else:
                    para.font.color.rgb = TEXT_SEC
            
            # Borders
            tc = cell._tc
            tcPr = tc.find(qn('a:tcPr'))
            if tcPr is None:
                tcPr = etree.SubElement(tc, qn('a:tcPr'))
            
            for border_name in ['lnL', 'lnR', 'lnT', 'lnB']:
                ln = etree.SubElement(tcPr, qn(f'a:{border_name}'), attrib={'w': '6350'})
                solidFill = etree.SubElement(ln, qn('a:solidFill'))
                srgbClr = etree.SubElement(solidFill, qn('a:srgbClr'), attrib={'val': '1E1648'})


def add_bullet_list(slide, left, top, width, height, items, font_size=13, color=TEXT_SEC, bullet_color=PURPLE):
    """Add a bulleted list."""
    txBox = add_textbox(slide, left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_after = Pt(5)
        p.space_before = Pt(2)
        p.alignment = PP_ALIGN.LEFT
        
        # Bullet character
        bullet_run = p.add_run()
        bullet_run.text = "▸  "
        bullet_run.font.size = Pt(font_size)
        bullet_run.font.color.rgb = bullet_color
        bullet_run.font.name = 'Calibri'
        bullet_run.font.bold = True
        
        # Text - handle highlight markers
        if '**' in item:
            parts = item.split('**')
            for j, part in enumerate(parts):
                r = p.add_run()
                r.text = part
                r.font.size = Pt(font_size)
                r.font.name = 'Calibri'
                if j % 2 == 1:
                    r.font.color.rgb = PURPLE_LIGHT
                    r.font.bold = True
                else:
                    r.font.color.rgb = color
        else:
            r = p.add_run()
            r.text = item
            r.font.size = Pt(font_size)
            r.font.color.rgb = color
            r.font.name = 'Calibri'
    
    return txBox


def add_flow_arrow(slide, left, top):
    """Add a flow arrow shape."""
    arrow = slide.shapes.add_shape(
        MSO_SHAPE.RIGHT_ARROW, left, top, Inches(0.4), Inches(0.3)
    )
    arrow.fill.solid()
    arrow.fill.fore_color.rgb = PURPLE
    arrow.line.fill.background()
    arrow.shadow.inherit = False
    return arrow


def add_flow_box(slide, left, top, width, height, icon, label, border=None):
    """Add a flow diagram box."""
    bc = border or RGBColor(0x2A, 0x20, 0x55)
    bg = add_shape_bg(slide, left, top, width, height, BG_CARD, bc)
    
    txBox = add_textbox(slide, left, top + Inches(0.05), width, height - Inches(0.05))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    
    r = p.add_run()
    r.text = icon
    r.font.size = Pt(18)
    r.font.name = 'Segoe UI Emoji'
    
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(2)
    r2 = p2.add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = WHITE
    r2.font.bold = True
    r2.font.name = 'Calibri'


def add_layer_row(slide, top, num, title, desc):
    """Add a numbered security layer row."""
    # Number circle
    circle = slide.shapes.add_shape(
        MSO_SHAPE.OVAL, MARGIN_LEFT + Inches(0.15), top + Inches(0.06), Inches(0.38), Inches(0.38)
    )
    circle.fill.solid()
    circle.fill.fore_color.rgb = PURPLE
    circle.line.fill.background()
    circle.shadow.inherit = False
    
    num_box = add_textbox(slide, MARGIN_LEFT + Inches(0.15), top + Inches(0.06), Inches(0.38), Inches(0.38))
    num_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = num_box.text_frame.paragraphs[0].add_run()
    r.text = str(num)
    r.font.size = Pt(14)
    r.font.color.rgb = WHITE
    r.font.bold = True
    r.font.name = 'Calibri'
    
    # Bar background
    bar = add_shape_bg(
        slide, MARGIN_LEFT + Inches(0.65), top,
        Inches(11.1), Inches(0.5), BG_CARD, RGBColor(0x2A, 0x20, 0x55)
    )
    
    # Left accent line on bar
    accent = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        MARGIN_LEFT + Inches(0.65), top, Pt(4), Inches(0.5)
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = PURPLE
    accent.line.fill.background()
    accent.shadow.inherit = False
    
    # Text
    txt = add_textbox(slide, MARGIN_LEFT + Inches(0.85), top + Inches(0.05), Inches(10.8), Inches(0.4))
    tf = txt.text_frame
    p = tf.paragraphs[0]
    
    r1 = p.add_run()
    r1.text = f"{title}:  "
    r1.font.size = Pt(13)
    r1.font.color.rgb = WHITE
    r1.font.bold = True
    r1.font.name = 'Calibri'
    
    r2 = p.add_run()
    r2.text = desc
    r2.font.size = Pt(12)
    r2.font.color.rgb = TEXT_SEC
    r2.font.name = 'Calibri'


def add_stat_card(slide, left, top, width, height, value, label):
    """Add a stat card with large number and label."""
    bg = add_shape_bg(slide, left, top, width, height, BG_CARD, RGBColor(0x2A, 0x20, 0x55))
    
    val_box = add_textbox(slide, left, top + Inches(0.15), width, Inches(0.6))
    val_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r = val_box.text_frame.paragraphs[0].add_run()
    r.text = value
    r.font.size = Pt(28)
    r.font.color.rgb = PURPLE_LIGHT
    r.font.bold = True
    r.font.name = 'Calibri'
    
    lbl_box = add_textbox(slide, left, top + Inches(0.7), width, Inches(0.4))
    lbl_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    r2 = lbl_box.text_frame.paragraphs[0].add_run()
    r2.text = label
    r2.font.size = Pt(10)
    r2.font.color.rgb = TEXT_SEC
    r2.font.name = 'Calibri'


# ═══════════════════════════════════════════════════════════════
# CREATE PRESENTATION
# ═══════════════════════════════════════════════════════════════
prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H

blank_layout = prs.slide_layouts[6]  # blank layout


# ═══════════════════════════════════════════════════════════════
# SLIDE 1: TITLE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)

# Large icon
icon_box = add_textbox(slide, Inches(5.4), Inches(0.8), Inches(2.5), Inches(1.2))
icon_box.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
r = icon_box.text_frame.paragraphs[0].add_run()
r.text = "🔐"
r.font.size = Pt(72)
r.font.name = 'Segoe UI Emoji'

# Title
title_box = add_textbox(slide, Inches(1.5), Inches(2.1), Inches(10.3), Inches(1.2))
tf = title_box.text_frame
tf.word_wrap = True
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Secure Cloud Storage System"
r.font.size = Pt(42)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = 'Calibri'

# Subtitle
sub_box = add_textbox(slide, Inches(2), Inches(3.2), Inches(9.3), Inches(0.5))
tf = sub_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "AES Encryption  •  Integrity Verification  •  Role-Based Access Control"
r.font.size = Pt(16)
r.font.color.rgb = PURPLE_LIGHT
r.font.name = 'Calibri'

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(3.85), Inches(2.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = PURPLE
line.line.fill.background()
line.shadow.inherit = False

# Authors
auth_box = add_textbox(slide, Inches(2.5), Inches(4.2), Inches(8.3), Inches(1.5))
tf = auth_box.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(6)
r = p.add_run()
r.text = "Ali Asjad Awan — FA23-BCS-028"
r.font.size = Pt(16)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = 'Calibri'

add_text_run(tf, "Muhammad Sohaib Liaqat — FA23-BCS-130", size=16, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, spacing_after=Pt(20))
add_text_run(tf, "Department of Computer Science", size=13, color=TEXT_SEC, alignment=PP_ALIGN.CENTER, spacing_after=Pt(2))
add_text_run(tf, "COMSATS University Islamabad  •  Spring 2026", size=13, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SLIDE 2: AGENDA
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Agenda")

left_items = [
    "Problem Statement & Objectives",
    "Literature Review",
    "System Architecture",
    "Methodology & Security Layers",
    "Technology Stack",
]
right_items = [
    "Implementation Deep Dive",
    "Encryption Workflow",
    "Authentication & Access Control",
    "Testing & Security Analysis",
    "Conclusion & Future Work",
]

add_bullet_list(slide, MARGIN_LEFT, CONTENT_TOP, Inches(5.5), Inches(4.5), left_items, font_size=14)
add_bullet_list(slide, Inches(6.5), CONTENT_TOP, Inches(5.5), Inches(4.5), right_items, font_size=14)


# ═══════════════════════════════════════════════════════════════
# SLIDE 3: PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Problem Statement")

cards_data = [
    ("💰", "$4.88M", "Average cost of a data breach in 2024\n(IBM Security Report)"),
    ("🔓", "Plaintext Storage", "Most cloud providers store data with\nserver-side keys they control"),
    ("⚠️", "No Integrity Check", "Files can be tampered with during\nstorage or transit undetected"),
]

card_w = Inches(3.7)
card_h = Inches(2.0)
gap = Inches(0.35)
start_x = MARGIN_LEFT + Inches(0.15)
for i, (icon, title, desc) in enumerate(cards_data):
    add_card(slide, start_x + i * (card_w + gap), CONTENT_TOP + Inches(0.3), card_w, card_h, icon, title, desc)

# Bottom callout
callout = add_textbox(slide, Inches(1.5), Inches(5.3), Inches(10.3), Inches(0.5))
tf = callout.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "→  We need client-controlled encryption + integrity verification + access control"
r.font.size = Pt(14)
r.font.color.rgb = PURPLE_LIGHT
r.font.bold = True
r.font.name = 'Calibri'


# ═══════════════════════════════════════════════════════════════
# SLIDE 4: RESEARCH OBJECTIVES
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Research Objectives")

objectives = [
    "Implement **AES encryption** (via Fernet) for files before cloud upload",
    "Implement **SHA-256 integrity** verification on every file download",
    "Build **secure authentication** with Bcrypt, rate limiting, and account lockout",
    "Implement **role-based access control** (RBAC) for user authorization",
    "Create **HMAC-signed tokens** for time-limited secure download links",
    "Implement **comprehensive audit logging** for security monitoring",
    "Deploy on **production cloud** infrastructure with HTTPS enforcement",
]
add_bullet_list(slide, MARGIN_LEFT, CONTENT_TOP, CONTENT_W, Inches(5.0), objectives, font_size=14)


# ═══════════════════════════════════════════════════════════════
# SLIDE 5: LITERATURE REVIEW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Literature Review — Comparison")

rows_data = [
    ["Feature", "Tresorit", "Boxcryptor", "SpiderOak", "Nextcloud", "Our System"],
    ["Client-Side Encryption", "✓", "✓", "✓", "Plugin", "✓"],
    ["Integrity Verification", "✓", "✗", "✓", "✗", "✓ SHA-256"],
    ["Time-Limited Links", "✓", "✗", "✗", "✓", "✓ HMAC"],
    ["RBAC", "✓", "✓", "✗", "✓", "✓"],
    ["Audit Logging", "✓", "✗", "✗", "✓", "✓"],
    ["CSRF Protection", "✓", "N/A", "N/A", "✓", "✓"],
    ["Open Source", "✗", "✗", "✗", "✓", "✓"],
]

table_w = Inches(11.5)
col_widths = [Inches(2.8), Inches(1.5), Inches(1.6), Inches(1.5), Inches(1.6), Inches(2.0)]
table = slide.shapes.add_table(len(rows_data), 6, MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.2), table_w, Inches(4.2)).table

for c, w in enumerate(col_widths):
    table.columns[c].width = w

for r_idx, row_data in enumerate(rows_data):
    for c_idx, val in enumerate(row_data):
        table.cell(r_idx, c_idx).text = val

style_table(table)

# Highlight "Our System" column
for r_idx in range(1, len(rows_data)):
    cell = table.cell(r_idx, 5)
    for para in cell.text_frame.paragraphs:
        para.font.color.rgb = GREEN
        para.font.bold = True


# ═══════════════════════════════════════════════════════════════
# SLIDE 6: SYSTEM ARCHITECTURE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "System Architecture")

flow_data = [
    ("🌐", "Client\nBrowser"),
    ("🔒", "HTTPS\nTLS"),
    ("⚙️", "Flask\nServer"),
    ("🔐", "Fernet\nEngine"),
    ("☁️", "Cloudinary\nCloud"),
]

box_w = Inches(1.6)
box_h = Inches(1.1)
arrow_w = Inches(0.5)
total_w = 5 * box_w + 4 * arrow_w
start_x = (SLIDE_W - total_w) / 2

flow_y = CONTENT_TOP + Inches(0.5)
for i, (icon, label) in enumerate(flow_data):
    x = start_x + i * (box_w + arrow_w)
    bc = RGBColor(0x2A, 0x20, 0x55)
    if i == 1:
        bc = RGBColor(0x16, 0x6D, 0x3B)
    elif i == 3:
        bc = RGBColor(0x7E, 0x22, 0xCE)
    add_flow_box(slide, x, flow_y, box_w, box_h, icon, label, bc)
    if i < 4:
        add_flow_arrow(slide, x + box_w + Inches(0.05), flow_y + Inches(0.4))

# DB connection row
db_y = flow_y + Inches(1.6)
add_flow_box(slide, start_x + 2 * (box_w + arrow_w), db_y, box_w, box_h, "⚙️", "Flask\nServer", RGBColor(0x2A, 0x20, 0x55))
add_flow_arrow(slide, start_x + 2 * (box_w + arrow_w) + box_w + Inches(0.05), db_y + Inches(0.4))
add_flow_box(slide, start_x + 3 * (box_w + arrow_w), db_y, box_w, box_h, "🗄️", "PostgreSQL\nDatabase", RGBColor(0x2A, 0x20, 0x55))

# Note
note_box = add_textbox(slide, Inches(1.5), Inches(5.5), Inches(10.3), Inches(0.5))
tf = note_box.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Files are encrypted in memory before leaving the server — Cloudinary only stores encrypted blobs"
r.font.size = Pt(13)
r.font.color.rgb = PURPLE_LIGHT
r.font.italic = True
r.font.name = 'Calibri'


# ═══════════════════════════════════════════════════════════════
# SLIDE 7: SECURITY LAYERS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "7 Security Layers")

layers = [
    ("Transport", "HTTPS enforcement via HSTS (31,536,000s max-age)"),
    ("Authentication", "Bcrypt hashing + secure session cookies"),
    ("Authorization", "RBAC with user/admin roles + ownership checks"),
    ("Input Validation", "CSRF tokens + magic-byte file type detection"),
    ("Encryption", "Fernet (AES-128-CBC + HMAC-SHA256) before upload"),
    ("Integrity", "SHA-256 hash verification on every download"),
    ("Monitoring", "Comprehensive audit logging with timestamps"),
]

for i, (title, desc) in enumerate(layers):
    y = CONTENT_TOP + Inches(0.1) + i * Inches(0.66)
    add_layer_row(slide, y, i + 1, title, desc)


# ═══════════════════════════════════════════════════════════════
# SLIDE 8: TECH STACK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Technology Stack")

tech = [
    ("🐍", "Python + Flask", "Backend framework"),
    ("🔐", "Fernet / AES", "File encryption"),
    ("🗄️", "PostgreSQL", "Database"),
    ("☁️", "Cloudinary", "Cloud storage"),
    ("#️⃣", "SHA-256", "Integrity hashing"),
    ("🔑", "Bcrypt", "Password hashing"),
    ("🛡️", "Flask-Talisman", "CSP + HSTS headers"),
    ("🚀", "Render", "Cloud deployment"),
]

card_w = Inches(2.7)
card_h = Inches(1.65)
gap_x = Inches(0.3)
gap_y = Inches(0.25)
cols = 4

for i, (icon, title, desc) in enumerate(tech):
    col = i % cols
    row = i // cols
    x = MARGIN_LEFT + Inches(0.1) + col * (card_w + gap_x)
    y = CONTENT_TOP + Inches(0.15) + row * (card_h + gap_y)
    add_card(slide, x, y, card_w, card_h, icon, title, desc)


# ═══════════════════════════════════════════════════════════════
# SLIDE 9: ENCRYPTION WORKFLOW
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Encryption Workflow")

# Left column - Upload
upload_title = add_textbox(slide, MARGIN_LEFT, CONTENT_TOP, Inches(5.5), Inches(0.4))
set_first_para(upload_title.text_frame, "📤  Upload (Encrypt)", size=18, color=PURPLE_LIGHT, bold=True)

upload_steps = [
    "User selects file for upload",
    "**Magic-byte** MIME type validation",
    "**SHA-256** hash computed and stored in DB",
    "**Fernet.encrypt()** — AES-128-CBC + HMAC",
    "Encrypted blob uploaded to **Cloudinary**",
    "Metadata persisted in **PostgreSQL**",
]
add_bullet_list(slide, MARGIN_LEFT, CONTENT_TOP + Inches(0.5), Inches(5.5), Inches(4.5), upload_steps, font_size=13)

# Right column - Download
dl_title = add_textbox(slide, Inches(6.8), CONTENT_TOP, Inches(5.5), Inches(0.4))
set_first_para(dl_title.text_frame, "📥  Download (Decrypt)", size=18, color=PURPLE_LIGHT, bold=True)

download_steps = [
    "**HMAC token** generated with 5-min expiry",
    "Token **signature + expiry** verified",
    "**Ownership / share** permission check",
    "Encrypted blob fetched from **Cloudinary**",
    "**Fernet.decrypt()** — verify HMAC tag",
    "**SHA-256** hash compared → serve file",
]
add_bullet_list(slide, Inches(6.8), CONTENT_TOP + Inches(0.5), Inches(5.5), Inches(4.5), download_steps, font_size=13)

# Divider line
div = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.5), CONTENT_TOP, Pt(2), Inches(5.0))
div.fill.solid()
div.fill.fore_color.rgb = RGBColor(0x2A, 0x20, 0x55)
div.line.fill.background()
div.shadow.inherit = False


# ═══════════════════════════════════════════════════════════════
# SLIDE 10: ENCRYPTION CODE
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Encryption — Code Implementation")

code_bg = add_shape_bg(slide, MARGIN_LEFT, CONTENT_TOP + Inches(0.1), CONTENT_W, Inches(5.0), RGBColor(0x05, 0x07, 0x15), RGBColor(0x2A, 0x20, 0x55))

code_text = """# Fernet encryption (AES-128-CBC + HMAC-SHA256)
from cryptography.fernet import Fernet
import hashlib

encryption_key = os.getenv("FERNET_KEY").encode()
cipher = Fernet(encryption_key)

# ── UPLOAD: Encrypt before cloud storage ──
data = file.read()
file_hash = hashlib.sha256(data).hexdigest()      # Integrity hash
encrypted_data = cipher.encrypt(data)              # AES encrypt
cloudinary.uploader.upload(
    BytesIO(encrypted_data), resource_type='raw'
)

# ── DOWNLOAD: Decrypt and verify integrity ──
encrypted_data = requests.get(file.storage_url).content
decrypted = cipher.decrypt(encrypted_data)         # AES decrypt + HMAC verify
current_hash = hashlib.sha256(decrypted).hexdigest()

if current_hash != file.file_hash:                 # Integrity check
    return "File integrity compromised!" """

code_box = add_textbox(slide, MARGIN_LEFT + Inches(0.3), CONTENT_TOP + Inches(0.25), Inches(11.3), Inches(4.7))
tf = code_box.text_frame
tf.word_wrap = True
set_first_para(tf, code_text, size=13, color=RGBColor(0xC0, 0xE0, 0xC0), font_name='Consolas')


# ═══════════════════════════════════════════════════════════════
# SLIDE 11: AUTH & ACCESS CONTROL
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Authentication & Access Control")

# Left card
add_shape_bg(slide, MARGIN_LEFT, CONTENT_TOP + Inches(0.1), Inches(5.6), Inches(4.8), BG_CARD, RGBColor(0x2A, 0x20, 0x55))
auth_title = add_textbox(slide, MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.2), Inches(5.2), Inches(0.4))
set_first_para(auth_title.text_frame, "🔑  Authentication", size=17, color=PURPLE_LIGHT, bold=True)

auth_items = [
    "**Bcrypt** password hashing with auto-generated salt",
    "**Rate limiting:** 5 requests/min on login endpoint",
    "**Account lockout** after 5 consecutive failed attempts",
    "**Password policy:** 8+ chars, upper, lower, digit, special",
    "**Secure cookies:** HttpOnly + Secure + SameSite=Lax",
]
add_bullet_list(slide, MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.7), Inches(5.2), Inches(4.0), auth_items, font_size=12)

# Right card
add_shape_bg(slide, Inches(6.6), CONTENT_TOP + Inches(0.1), Inches(5.6), Inches(4.8), BG_CARD, RGBColor(0x2A, 0x20, 0x55))
rbac_title = add_textbox(slide, Inches(6.8), CONTENT_TOP + Inches(0.2), Inches(5.2), Inches(0.4))
set_first_para(rbac_title.text_frame, "👥  Authorization (RBAC)", size=17, color=PURPLE_LIGHT, bold=True)

rbac_items = [
    "**User** role — Manage own files, share, analytics",
    "**Admin** role — View all files, users, audit logs",
    "**Ownership checks** on every file operation",
    "**Share permissions:** view / download with 7-day expiry",
    "**Instant revocation** terminates access immediately",
]
add_bullet_list(slide, Inches(6.8), CONTENT_TOP + Inches(0.7), Inches(5.2), Inches(4.0), rbac_items, font_size=12)


# ═══════════════════════════════════════════════════════════════
# SLIDE 12: FILE SHARING SYSTEM
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "File Sharing System")

# Flow diagram
flow_items = [
    ("👤", "Owner"),
    ("📝", "Share Form"),
    ("🔔", "Notification"),
    ("👤", "Recipient"),
]
fbox_w = Inches(1.8)
fbox_h = Inches(0.95)
farrow_w = Inches(0.45)
total = len(flow_items) * fbox_w + (len(flow_items) - 1) * farrow_w
fstart = (SLIDE_W - total) / 2
fy = CONTENT_TOP + Inches(0.2)

for i, (icon, label) in enumerate(flow_items):
    x = fstart + i * (fbox_w + farrow_w)
    add_flow_box(slide, x, fy, fbox_w, fbox_h, icon, label)
    if i < len(flow_items) - 1:
        add_flow_arrow(slide, x + fbox_w + Inches(0.03), fy + Inches(0.32))

# Two cards below
card_y = CONTENT_TOP + Inches(1.6)
add_shape_bg(slide, MARGIN_LEFT, card_y, Inches(5.6), Inches(3.7), BG_CARD, RGBColor(0x2A, 0x20, 0x55))
ft1 = add_textbox(slide, MARGIN_LEFT + Inches(0.2), card_y + Inches(0.15), Inches(5.2), Inches(0.35))
set_first_para(ft1.text_frame, "Features", size=16, color=PURPLE_LIGHT, bold=True)
add_bullet_list(slide, MARGIN_LEFT + Inches(0.2), card_y + Inches(0.55), Inches(5.2), Inches(3.0), [
    "View / Download permission levels",
    "7-day automatic share expiry",
    "Download counter tracking per share",
    "Real-time notifications for recipients",
    "Idempotent share creation/updates",
], font_size=12)

add_shape_bg(slide, Inches(6.6), card_y, Inches(5.6), Inches(3.7), BG_CARD, RGBColor(0x2A, 0x20, 0x55))
ft2 = add_textbox(slide, Inches(6.8), card_y + Inches(0.15), Inches(5.2), Inches(0.35))
set_first_para(ft2.text_frame, "Security Controls", size=16, color=PURPLE_LIGHT, bold=True)
add_bullet_list(slide, Inches(6.8), card_y + Inches(0.55), Inches(5.2), Inches(3.0), [
    "Ownership verification before sharing",
    "Cannot share files with yourself",
    "Unique cryptographic share tokens",
    "Instant revocation of access",
    "Expired shares automatically blocked",
], font_size=12)


# ═══════════════════════════════════════════════════════════════
# SLIDE 13: SECURITY HARDENING
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Security Hardening")

hardening = [
    ("🛡️", "CSRF Protection", "Flask-WTF unique tokens per\nsession, verified server-side"),
    ("📋", "Content Security Policy", "Restricts scripts/styles to\ntrusted origins only"),
    ("🔒", "HSTS Enforcement", "Forces HTTPS with 1-year\nmax-age directive"),
    ("🚫", "Malware Blocking", "Magic-byte analysis + PE/ELF\nsignature detection"),
    ("⏱️", "Token Expiry", "HMAC-signed download tokens\nexpire after 5 minutes"),
    ("📊", "Audit Logging", "All actions logged with user\nID and timestamp"),
]

card_w = Inches(3.7)
card_h = Inches(1.95)
gap_x = Inches(0.35)
gap_y = Inches(0.25)

for i, (icon, title, desc) in enumerate(hardening):
    col = i % 3
    row = i // 3
    x = MARGIN_LEFT + Inches(0.1) + col * (card_w + gap_x)
    y = CONTENT_TOP + Inches(0.15) + row * (card_h + gap_y)
    add_card(slide, x, y, card_w, card_h, icon, title, desc)


# ═══════════════════════════════════════════════════════════════
# SLIDE 14: DATABASE DESIGN
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Database Design — 5 Models")

db_models = [
    ("👤", "User", "id, username, email, password\n(Bcrypt), role, failed_attempts"),
    ("📄", "File", "id, filename, encrypted_name,\ncloud_url, owner_id, file_hash"),
    ("🤝", "SharedFile", "file_id, shared_by, shared_with,\npermission, expiry, downloads"),
    ("📋", "Log", "id, user_id, action, timestamp\n— records all security events"),
    ("🔔", "Notification", "user_id, message, read,\ncreated_at, file_id, share_id"),
]

card_w = Inches(2.18)
card_h = Inches(2.0)
gap = Inches(0.25)

for i, (icon, title, desc) in enumerate(db_models):
    x = MARGIN_LEFT + Inches(0.15) + i * (card_w + gap)
    add_card(slide, x, CONTENT_TOP + Inches(0.3), card_w, card_h, icon, title, desc)


# ═══════════════════════════════════════════════════════════════
# SLIDE 15: APPLICATION SCREENSHOTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Application — Key Screens")

screens = [
    ("🏠", "Landing Page", "Modern glassmorphism UI with\nfeature showcase & CTA"),
    ("🔐", "Login Page", "CSRF-protected login with\nforgot password support"),
    ("📊", "Dashboard", "Quick actions, security\nfeatures, password change"),
    ("📤", "Upload Page", "File upload with magic-byte\ntype validation"),
    ("📁", "File Management", "Files list with download,\nshare, and delete actions"),
    ("🔔", "Notifications", "Real-time share alerts with\nmark-as-read feature"),
]

card_w = Inches(3.7)
card_h = Inches(1.95)
gap_x = Inches(0.35)
gap_y = Inches(0.25)

for i, (icon, title, desc) in enumerate(screens):
    col = i % 3
    row = i // 3
    x = MARGIN_LEFT + Inches(0.1) + col * (card_w + gap_x)
    y = CONTENT_TOP + Inches(0.15) + row * (card_h + gap_y)
    add_card(slide, x, y, card_w, card_h, icon, title, desc)


# ═══════════════════════════════════════════════════════════════
# SLIDE 16: SECURITY TESTING RESULTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Security Testing Results")

sec_rows = [
    ["Attack Type", "Method", "Expected", "Result"],
    ["CSRF Attack", "POST without token", "400 Rejected", "PASS ✓"],
    ["Brute Force", "10 rapid login attempts", "Rate limited + locked", "PASS ✓"],
    ["IDOR", "Access other user's file", "403 Forbidden", "PASS ✓"],
    ["Token Replay", "Reuse expired token", "Link expired", "PASS ✓"],
    ["Token Forgery", "Modified HMAC signature", "Invalid link", "PASS ✓"],
    ["Malware Upload", "Upload .exe (MZ header)", "Rejected", "PASS ✓"],
    ["SQL Injection", "SQL in form fields", "Parameterized query", "PASS ✓"],
    ["XSS", "Script tag in input", "HTML escaped", "PASS ✓"],
    ["Cookie Hijack", "JS cookie access attempt", "HttpOnly blocked", "PASS ✓"],
    ["Downgrade", "HTTP access attempt", "HSTS redirect", "PASS ✓"],
]

col_widths = [Inches(2.2), Inches(3.5), Inches(3.5), Inches(1.8)]
table = slide.shapes.add_table(len(sec_rows), 4, MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.1), Inches(11.0), Inches(5.3)).table

for c, w in enumerate(col_widths):
    table.columns[c].width = w

for r_idx, row_data in enumerate(sec_rows):
    for c_idx, val in enumerate(row_data):
        table.cell(r_idx, c_idx).text = val

style_table(table)

# Green highlight for PASS column
for r_idx in range(1, len(sec_rows)):
    cell = table.cell(r_idx, 3)
    for para in cell.text_frame.paragraphs:
        para.font.color.rgb = GREEN
        para.font.bold = True


# ═══════════════════════════════════════════════════════════════
# SLIDE 17: PERFORMANCE RESULTS
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Performance Results")

stats = [
    ("<5ms", "100KB Encryption"),
    ("~12ms", "1MB Encryption"),
    ("~45ms", "5MB Encryption"),
    ("~90ms", "10MB Encryption"),
]

stat_w = Inches(2.7)
stat_h = Inches(1.4)
gap = Inches(0.3)
sx = MARGIN_LEFT + Inches(0.2)

for i, (val, lbl) in enumerate(stats):
    x = sx + i * (stat_w + gap)
    add_stat_card(slide, x, CONTENT_TOP + Inches(0.5), stat_w, stat_h, val, lbl)

# Note
note = add_textbox(slide, Inches(1.5), Inches(4.0), Inches(10.3), Inches(0.5))
tf = note.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Fernet encryption adds negligible latency — network upload/download is the dominant factor"
r.font.size = Pt(14)
r.font.color.rgb = PURPLE_LIGHT
r.font.italic = True
r.font.name = 'Calibri'

# Performance table
perf_rows = [
    ["File Size", "Encryption Time", "Decryption Time", "Total Upload Time"],
    ["100 KB", "< 5 ms", "< 5 ms", "~1.2 s"],
    ["1 MB", "~12 ms", "~10 ms", "~2.5 s"],
    ["5 MB", "~45 ms", "~40 ms", "~5.8 s"],
    ["10 MB", "~90 ms", "~85 ms", "~9.2 s"],
]

perf_table = slide.shapes.add_table(len(perf_rows), 4, MARGIN_LEFT + Inches(1.5), Inches(4.6), Inches(9.0), Inches(2.3)).table
perf_col_w = [Inches(2.0), Inches(2.3), Inches(2.3), Inches(2.4)]
for c, w in enumerate(perf_col_w):
    perf_table.columns[c].width = w

for r_idx, row_data in enumerate(perf_rows):
    for c_idx, val in enumerate(row_data):
        perf_table.cell(r_idx, c_idx).text = val

style_table(perf_table)


# ═══════════════════════════════════════════════════════════════
# SLIDE 18: CIA TRIAD
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "CIA Triad Analysis")

cia = [
    ("🔒", "Confidentiality", "Fernet AES-128-CBC encryption before\ncloud upload. Bcrypt password hashing.\nHttpOnly + Secure session cookies."),
    ("✅", "Integrity", "SHA-256 hash verification on every\ndownload. Fernet HMAC authentication.\nCSRF token validation on POST."),
    ("🌐", "Availability", "Managed cloud deployment on Render.\nRate limiting prevents DoS attacks.\nAutomatic restart on failure."),
]

card_w = Inches(3.7)
card_h = Inches(2.6)
gap = Inches(0.35)

for i, (icon, title, desc) in enumerate(cia):
    x = MARGIN_LEFT + Inches(0.15) + i * (card_w + gap)
    add_card(slide, x, CONTENT_TOP + Inches(0.5), card_w, card_h, icon, title, desc)


# ═══════════════════════════════════════════════════════════════
# SLIDE 19: OWASP TOP 10
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "OWASP Top 10 Compliance")

owasp_rows = [
    ["OWASP Risk", "Status", "How We Address It"],
    ["A01: Broken Access Control", "Mitigated ✓", "RBAC + ownership checks + token expiry"],
    ["A02: Cryptographic Failures", "Mitigated ✓", "Fernet (AES) + Bcrypt + SHA-256"],
    ["A03: Injection", "Mitigated ✓", "SQLAlchemy ORM parameterized queries"],
    ["A04: Insecure Design", "Mitigated ✓", "Defense-in-depth + threat modeling"],
    ["A05: Security Misconfiguration", "Mitigated ✓", "CSP + HSTS + secure cookie flags"],
    ["A07: Auth Failures", "Mitigated ✓", "Rate limiting + lockout + strong passwords"],
    ["A08: Data Integrity Failures", "Mitigated ✓", "SHA-256 verification + HMAC tokens"],
    ["A09: Logging Failures", "Mitigated ✓", "Comprehensive audit logging system"],
]

col_widths = [Inches(3.8), Inches(2.0), Inches(5.5)]
table = slide.shapes.add_table(len(owasp_rows), 3, MARGIN_LEFT + Inches(0.15), CONTENT_TOP + Inches(0.15), Inches(11.3), Inches(5.0)).table

for c, w in enumerate(col_widths):
    table.columns[c].width = w

for r_idx, row_data in enumerate(owasp_rows):
    for c_idx, val in enumerate(row_data):
        table.cell(r_idx, c_idx).text = val

style_table(table)

for r_idx in range(1, len(owasp_rows)):
    cell = table.cell(r_idx, 1)
    for para in cell.text_frame.paragraphs:
        para.font.color.rgb = GREEN
        para.font.bold = True


# ═══════════════════════════════════════════════════════════════
# SLIDE 20: LIMITATIONS & FUTURE WORK
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Limitations & Future Work")

# Left - Limitations
add_shape_bg(slide, MARGIN_LEFT, CONTENT_TOP + Inches(0.1), Inches(5.6), Inches(4.8), BG_CARD, RGBColor(0x2A, 0x20, 0x55))
lim_title = add_textbox(slide, MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.2), Inches(5.2), Inches(0.4))
set_first_para(lim_title.text_frame, "⚠️  Current Limitations", size=17, color=YELLOW, bold=True)

add_bullet_list(slide, MARGIN_LEFT + Inches(0.2), CONTENT_TOP + Inches(0.7), Inches(5.2), Inches(4.0), [
    "Server has access to plaintext during encryption",
    "Single Fernet key for all users (shared key)",
    "No two-factor authentication (2FA)",
    "No file versioning support",
    "No client-side (browser) encryption",
], font_size=12)

# Right - Future
add_shape_bg(slide, Inches(6.6), CONTENT_TOP + Inches(0.1), Inches(5.6), Inches(4.8), BG_CARD, RGBColor(0x2A, 0x20, 0x55))
fut_title = add_textbox(slide, Inches(6.8), CONTENT_TOP + Inches(0.2), Inches(5.2), Inches(0.4))
set_first_para(fut_title.text_frame, "🚀  Future Enhancements", size=17, color=GREEN, bold=True)

add_bullet_list(slide, Inches(6.8), CONTENT_TOP + Inches(0.7), Inches(5.2), Inches(4.0), [
    "**End-to-End Encryption** — Web Crypto API",
    "**Two-Factor Auth** — TOTP via Google Authenticator",
    "**Key Rotation** — Automated periodic re-encryption",
    "**Zero-Knowledge** — Client-side key derivation",
    "**GDPR/HIPAA** — Compliance and audit features",
], font_size=12)


# ═══════════════════════════════════════════════════════════════
# SLIDE 21: CONCLUSION
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)
add_title_bar(slide, "Conclusion")

# Summary card
add_shape_bg(slide, MARGIN_LEFT, CONTENT_TOP + Inches(0.1), CONTENT_W, Inches(1.8), BG_CARD, RGBColor(0x2A, 0x20, 0x55))
summary = add_textbox(slide, MARGIN_LEFT + Inches(0.3), CONTENT_TOP + Inches(0.2), Inches(11.3), Inches(1.6))
tf = summary.text_frame
tf.word_wrap = True
set_first_para(tf, (
    "We designed and deployed a production-ready secure cloud storage system that encrypts files "
    "with AES-128-CBC (Fernet) before cloud upload, verifies integrity with SHA-256 on every download, "
    "protects credentials with Bcrypt hashing, and controls access through RBAC and HMAC-signed "
    "time-limited download tokens. All 10 security tests passed with zero vulnerabilities."
), size=13, color=TEXT_SEC, font_name='Calibri')

# Stats
stat_data = [
    ("10/10", "Security Tests\nPassed"),
    ("7", "OWASP Risks\nMitigated"),
    ("7", "Security\nLayers"),
    ("<90ms", "Max Encryption\nOverhead"),
]

stat_w = Inches(2.7)
stat_h = Inches(1.7)
gap = Inches(0.3)
sy = CONTENT_TOP + Inches(2.2)

for i, (val, lbl) in enumerate(stat_data):
    x = MARGIN_LEFT + Inches(0.2) + i * (stat_w + gap)
    add_stat_card(slide, x, sy, stat_w, stat_h, val, lbl)


# ═══════════════════════════════════════════════════════════════
# SLIDE 22: Q&A
# ═══════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
set_slide_bg(slide, BG_DARK)
add_morph_transition(slide)

# Large icon
qa_icon = add_textbox(slide, Inches(5.2), Inches(0.8), Inches(3), Inches(1.5))
qa_icon.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
r = qa_icon.text_frame.paragraphs[0].add_run()
r.text = "💬"
r.font.size = Pt(80)
r.font.name = 'Segoe UI Emoji'

# Q&A text
qa_title = add_textbox(slide, Inches(1.5), Inches(2.5), Inches(10.3), Inches(1.0))
tf = qa_title.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Questions & Discussion"
r.font.size = Pt(42)
r.font.color.rgb = PURPLE_LIGHT
r.font.bold = True
r.font.name = 'Calibri'

# Thank you
thanks = add_textbox(slide, Inches(2.5), Inches(3.6), Inches(8.3), Inches(0.5))
tf = thanks.text_frame
p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
r = p.add_run()
r.text = "Thank you for your attention!"
r.font.size = Pt(16)
r.font.color.rgb = TEXT_SEC
r.font.name = 'Calibri'

# Accent line
line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(4.3), Inches(2.3), Pt(3))
line.fill.solid()
line.fill.fore_color.rgb = PURPLE
line.line.fill.background()
line.shadow.inherit = False

# Authors
auth = add_textbox(slide, Inches(2.5), Inches(4.6), Inches(8.3), Inches(1.5))
tf = auth.text_frame
tf.word_wrap = True

p = tf.paragraphs[0]
p.alignment = PP_ALIGN.CENTER
p.space_after = Pt(4)
r = p.add_run()
r.text = "Ali Asjad Awan — FA23-BCS-028"
r.font.size = Pt(15)
r.font.color.rgb = WHITE
r.font.bold = True
r.font.name = 'Calibri'

add_text_run(tf, "Muhammad Sohaib Liaqat — FA23-BCS-130", size=15, color=WHITE, bold=True, alignment=PP_ALIGN.CENTER, spacing_after=Pt(16))
add_text_run(tf, "COMSATS University Islamabad  •  Spring 2026", size=12, color=TEXT_SEC, alignment=PP_ALIGN.CENTER)


# ═══════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), 'presentation.pptx')
prs.save(output_path)
print(f"Presentation saved to: {output_path}")
print(f"Total slides: {len(prs.slides)}")
print("Done!")
